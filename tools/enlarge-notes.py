#!/usr/bin/env python3
"""Make the speaker notes readable in the built PPTX and HTML.

PowerPoint renders speaker notes at 12pt. On a venue laptop, in Presenter View, at arm's
length, under stage light, that is not a size you can glance at — and glancing is the only
way notes get used. Nothing upstream can fix it: Marp writes the notes into the notes
placeholder and PowerPoint's own template decides the size, so the fix has to happen after
the file exists.

Optional by design. python-pptx is not needed to build the deck, only to make the notes
comfortable, so `make pptx` skips this step with a hint rather than failing when the module
is missing.

The HTML presenter view has the same problem and a much easier fix: its note pane scales
off a CSS variable, `--bespoke-marp-note-font-scale`, which the view exposes as a control
nobody finds. Setting a larger default costs one style rule. If a future Marp renames the
variable the rule stops doing anything, which is the right way for this to fail.

Usage: python3 tools/enlarge-notes.py build/deck.pptx [--pt 20]
       python3 tools/enlarge-notes.py build/deck.html [--scale 1.6]
"""

from __future__ import annotations

import argparse
import pathlib
import sys


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=pathlib.Path, help="build/deck.pptx or build/deck.html")
    ap.add_argument("--pt", type=float, default=20.0, help="PPTX point size (default 20)")
    ap.add_argument("--scale", type=float, default=1.6,
                    help="HTML presenter-view note scale (default 1.6)")
    args = ap.parse_args()

    if args.pptx.suffix == ".html":
        html = args.pptx.read_text(encoding="utf-8")
        rule = (f"<style>:root{{--bespoke-marp-note-font-scale:{args.scale:g}}}</style>")
        if rule in html:
            return 0
        args.pptx.write_text(html.replace("</head>", rule + "</head>", 1), encoding="utf-8")
        print(f"enlarge-notes.py: presenter-view notes scaled to {args.scale:g}x")
        return 0

    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ModuleNotFoundError:
        print("enlarge-notes.py: python-pptx not installed, leaving the notes at 12pt.\n"
              "                  pip install python-pptx  to get readable ones.",
              file=sys.stderr)
        return 0

    deck = Presentation(str(args.pptx))
    slides = 0
    for slide in deck.slides:
        if not slide.has_notes_slide:
            continue
        frame = slide.notes_slide.notes_text_frame
        if not frame.text.strip():
            continue
        for paragraph in frame.paragraphs:
            # A run carries its own size and overrides the paragraph, so set both: a
            # paragraph with no runs (a blank line) would otherwise keep the default.
            paragraph.font.size = Pt(args.pt)
            for run in paragraph.runs:
                run.font.size = Pt(args.pt)
        slides += 1

    deck.save(str(args.pptx))
    print(f"enlarge-notes.py: speaker notes set to {args.pt:g}pt on {slides} slide(s)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
